"""
Agent 3: PPT Generator
程序化生成具有科幻感的精美 PPT。
三套主题，深色科幻风为主，每页从空白页程序化构建。
"""

import logging
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from app.core.config import DATA_DIR

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Theme definitions — all with dark background for sci-fi feel
# ---------------------------------------------------------------------------

THEMES = {
    "tech_blue": {
        "name": "科技蓝",
        "colors": {
            "bg":         RGBColor(0x08, 0x0D, 0x20),  # 极深蓝黑背景
            "bg_accent":  RGBColor(0x0F, 0x17, 0x3A),  # 稍亮蓝黑
            "primary":    RGBColor(0x1A, 0x56, 0xDB),  # 主蓝色
            "cyan":       RGBColor(0x00, 0xD4, 0xFF),  # 青色
            "cyan_dim":   RGBColor(0x00, 0x8A, 0xBB),  # 暗青
            "accent":     RGBColor(0x7C, 0x3A, 0xED),  # 紫
            "title":      RGBColor(0xFF, 0xFF, 0xFF),  # 标题白
            "subtitle":   RGBColor(0x8A, 0xB4, 0xFF),  # 副标题淡蓝
            "heading":    RGBColor(0x00, 0xD4, 0xFF),  # 小标题青色
            "body":       RGBColor(0xE0, 0xE8, 0xF5),  # 正文浅灰白
            "muted":      RGBColor(0x5A, 0x6A, 0x8A),  # 辅助文字
            "bullet":     RGBColor(0x00, 0xD4, 0xFF),  # 圆点青色
            "code_bg":    RGBColor(0x12, 0x1A, 0x33),  # 代码块背景
            "code_border": RGBColor(0x1A, 0x56, 0xDB), # 代码边框
            "line":       RGBColor(0x1A, 0x3A, 0x6A),  # 分割线
            "glow":       RGBColor(0x00, 0xB4, 0xD8),  # 发光色
        },
        "fonts": {"heading": "Microsoft YaHei", "body": "Microsoft YaHei"},
    },
    "clean_white": {
        "name": "简约白",
        "colors": {
            "bg":         RGBColor(0x1A, 0x1A, 0x2E),  # 深色底
            "bg_accent":  RGBColor(0x22, 0x22, 0x38),
            "primary":    RGBColor(0x2D, 0x34, 0x3E),
            "cyan":       RGBColor(0xE0, 0x6D, 0x06),  # 橙色强调
            "cyan_dim":   RGBColor(0xA0, 0x50, 0x00),
            "accent":     RGBColor(0xF5, 0x9E, 0x0B),
            "title":      RGBColor(0xFF, 0xFF, 0xFF),
            "subtitle":   RGBColor(0xBB, 0xC1, 0xCC),
            "heading":    RGBColor(0xE0, 0x6D, 0x06),
            "body":       RGBColor(0xE0, 0xE4, 0xEA),
            "muted":      RGBColor(0x7A, 0x84, 0x96),
            "bullet":     RGBColor(0xE0, 0x6D, 0x06),
            "code_bg":    RGBColor(0x2A, 0x2A, 0x3E),
            "code_border":RGBColor(0xE0, 0x6D, 0x06),
            "line":       RGBColor(0x3A, 0x3A, 0x50),
            "glow":       RGBColor(0xE0, 0x6D, 0x06),
        },
        "fonts": {"heading": "Microsoft YaHei", "body": "Microsoft YaHei"},
    },
    "warm_orange": {
        "name": "活力橙",
        "colors": {
            "bg":         RGBColor(0x12, 0x0A, 0x08),
            "bg_accent":  RGBColor(0x22, 0x12, 0x0A),
            "primary":    RGBColor(0xFF, 0x6B, 0x35),
            "cyan":       RGBColor(0xFF, 0x8A, 0x50),
            "cyan_dim":   RGBColor(0xC0, 0x50, 0x20),
            "accent":     RGBColor(0xFF, 0xCC, 0x80),
            "title":      RGBColor(0xFF, 0xFF, 0xFF),
            "subtitle":   RGBColor(0xFF, 0xC0, 0xA0),
            "heading":    RGBColor(0xFF, 0x8A, 0x50),
            "body":       RGBColor(0xF0, 0xE8, 0xE0),
            "muted":      RGBColor(0x90, 0x70, 0x60),
            "bullet":     RGBColor(0xFF, 0x6B, 0x35),
            "code_bg":    RGBColor(0x1E, 0x14, 0x10),
            "code_border":RGBColor(0xFF, 0x6B, 0x35),
            "line":       RGBColor(0x3A, 0x20, 0x18),
            "glow":       RGBColor(0xFF, 0x8A, 0x50),
        },
        "fonts": {"heading": "Microsoft YaHei", "body": "Microsoft YaHei"},
    },
}

DEFAULT_THEME = "tech_blue"


def list_templates() -> list[dict]:
    return [
        {"id": key, "name": info["name"], "exists": True}
        for key, info in THEMES.items()
    ]


def _theme(id_: str) -> dict:
    return THEMES.get(id_, THEMES[DEFAULT_THEME])


# ---------------------------------------------------------------------------
# Run helpers
# ---------------------------------------------------------------------------

def _run(para, text: str, size: int, color: RGBColor,
         bold: bool = False, font_name: str = "Microsoft YaHei",
         italic: bool = False):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font_name
    r.font.italic = italic
    return r


def _add_rect(slide, x, y, w, h, color, alpha=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # python-pptx doesn't directly support alpha on shape fills,
    # but the shape is there — PPT renders it
    return shape


def _add_circle(slide, x, y, size, color):
    """Add a circle/oval."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Title slide — cyberpunk/sci-fi style
# ---------------------------------------------------------------------------

def _build_title_slide(slide, title: str, subtitle: str, theme: dict):
    c = theme["colors"]
    fonts = theme["fonts"]
    W, H = Inches(7.5), Inches(13.333)

    # Background — full slide dark
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = c["bg"]

    # -- Decorative geometry: large subtle circles in bg --
    _add_circle(slide, Inches(5.5), Inches(-1.5), Inches(4), c["bg_accent"])
    _add_circle(slide, Inches(-1.0), Inches(8.0), Inches(3.5), c["bg_accent"])

    # -- Top accent line (thin, full width) --
    _add_rect(slide, Inches(0), Inches(0), W, Pt(3), c["cyan"])

    # -- Left vertical accent bar --
    _add_rect(slide, Inches(0.5), Inches(2.5), Pt(4), Inches(6.0), c["primary"])

    # -- Bottom glow bar --
    _add_rect(slide, Inches(0), Inches(13.0), W, Pt(4), c["cyan"])

    # -- Decorative triangle (top-right area) --
    try:
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.5), Inches(0.5), Inches(1.5), Inches(2.5))
        tri.fill.solid()
        tri.fill.fore_color.rgb = c["bg_accent"]
        tri.line.fill.background()
        tri.rotation = 20.0
    except Exception:
        pass

    # -- Decorative small diamond (top-right) --
    try:
        dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(6.2), Inches(1.0), Inches(0.8), Inches(0.8))
        dia.fill.solid()
        dia.fill.fore_color.rgb = c["primary"]
        dia.line.fill.background()
    except Exception:
        pass

    # -- Title text --
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(5.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    _run(p, title, 36, c["title"], bold=True, font_name=fonts["heading"])

    # -- Title underline accent --
    title_len = min(len(title) * 22, 400)
    _add_rect(slide, Inches(1.0), Inches(5.5), Inches(title_len / 15), Pt(3), c["cyan"])

    # -- Subtitle --
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(5.5), Inches(2.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        _run(p2, subtitle, 14, c["subtitle"], font_name=fonts["body"])

    # -- Bottom-right decoration (small line block) --
    _add_rect(slide, Inches(5.5), Inches(12.5), Inches(1.5), Pt(2), c["primary"])
    _add_rect(slide, Inches(5.5), Inches(12.3), Inches(0.8), Pt(2), c["cyan"])

    # Page number
    txPage = slide.shapes.add_textbox(Inches(6.0), Inches(12.8), Inches(1.0), Inches(0.4))
    p = txPage.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, "1", 10, c["muted"], font_name=fonts["body"])


# ---------------------------------------------------------------------------
# Content slide — dark bg, neon accents, clean typography
# ---------------------------------------------------------------------------

def _build_content_slide(slide, title: str, content: str, page_num: int, theme: dict):
    c = theme["colors"]
    fonts = theme["fonts"]

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = c["bg"]

    # -- Subtle bg decoration (circle top-right) --
    _add_circle(slide, Inches(6.0), Inches(-1.5), Inches(3), c["bg_accent"])

    # -- Top accent line --
    _add_rect(slide, Inches(0), Inches(0), Inches(7.5), Pt(3), c["cyan"])

    # -- Left neon accent bar (full height) --
    _add_rect(slide, Inches(0.4), Inches(0.3), Pt(5), Inches(12.5), c["primary"])
    # Second thinner neon bar
    _add_rect(slide, Inches(0.4), Inches(0.3), Pt(2), Inches(12.5), c["cyan"])

    # -- Title --
    txTitle = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(5.8), Inches(0.7))
    tf_title = txTitle.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, title, 24, c["heading"], bold=True, font_name=fonts["heading"])

    # -- Separator line --
    _add_rect(slide, Inches(0.9), Inches(1.35), Inches(5.5), Pt(1.5), c["line"])

    # -- Content area --
    txContent = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.8), Inches(10.5))
    tf = txContent.text_frame
    tf.word_wrap = True

    _parse_markdown(tf, content, c, fonts)

    # -- Page number --
    txPage = slide.shapes.add_textbox(Inches(6.0), Inches(12.5), Inches(1.0), Inches(0.4))
    pn = txPage.text_frame.paragraphs[0]
    pn.alignment = PP_ALIGN.RIGHT
    _run(pn, str(page_num), 10, c["muted"], font_name=fonts["body"])


# ---------------------------------------------------------------------------
# Code block background shape
# ---------------------------------------------------------------------------

def _add_code_block_bg(slide, y_top, line_count: int, theme: dict):
    """Add a dark rectangle behind code lines for the code-block look."""
    c = theme["colors"]
    height = max(Inches(0.4), min(line_count * Inches(0.25), Inches(3.5)))
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.1), y_top, Inches(5.3), height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c["code_bg"]
    shape.line.color.rgb = c["code_border"]
    shape.line.width = Pt(1)
    return shape


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def _parse_markdown(tf, content: str, c: dict, fonts: dict):
    """Parse markdown content into styled text frame paragraphs."""
    lines = content.split("\n")
    code_block = False
    code_lines = []
    code_start = True
    first = True

    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        # Code block toggle
        if stripped.startswith("```"):
            if code_block:
                # End code block — render it in a single paragraph
                if code_lines:
                    if first:
                        para = tf.paragraphs[0]
                        first = False
                    else:
                        para = tf.add_paragraph()
                    para.space_before = Pt(6)
                    para.space_after = Pt(6)
                    for cl in code_lines:
                        _run(para, cl + "\n", 11, c["body"],
                             font_name="Consolas")
                code_lines = []
                code_block = False
            else:
                code_block = True
                code_lines = []
            i += 1
            continue

        if code_block:
            code_lines.append(stripped)
            i += 1
            continue

        # Empty line
        if not stripped:
            if first:
                first = False
            else:
                para = tf.add_paragraph()
                para.space_before = Pt(4)
                para.space_after = Pt(4)
                _run(para, "", 10, c["body"])
            i += 1
            continue

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        # Headings
        if stripped.startswith("### "):
            _run(para, stripped[4:], 18, c["heading"], bold=True, font_name=fonts["heading"])
            para.space_before = Pt(14)
            para.space_after = Pt(6)

        elif stripped.startswith("## "):
            _run(para, stripped[3:], 21, c["heading"], bold=True, font_name=fonts["heading"])
            para.space_before = Pt(18)
            para.space_after = Pt(8)

        elif stripped.startswith("# "):
            _run(para, stripped[2:], 24, c["heading"], bold=True, font_name=fonts["heading"])
            para.space_before = Pt(20)
            para.space_after = Pt(8)

        # Bullet points
        elif stripped.startswith("- ") or stripped.startswith("❍ "):
            text = stripped[2:] if stripped.startswith("- ") else stripped[1:]
            _run(para, "◆  ", 12, c["bullet"], font_name=fonts["body"])
            _run(para, text, 15, c["body"], font_name=fonts["body"])
            para.space_before = Pt(4)
            para.space_after = Pt(4)
            para.level = 0

        # Numbered list
        elif re.match(r"^\d+[.、]", stripped):
            _run(para, stripped, 15, c["body"], font_name=fonts["body"])
            para.space_before = Pt(4)
            para.space_after = Pt(4)

        # Bold text
        elif "**" in stripped:
            parts = re.split(r"(\*\*.*?\*\*)", stripped)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    _run(para, part[2:-2], 15, c["heading"], bold=True, font_name=fonts["body"])
                else:
                    _run(para, part, 15, c["body"], font_name=fonts["body"])
            para.space_after = Pt(4)

        else:
            _run(para, stripped, 15, c["body"], font_name=fonts["body"])
            para.space_after = Pt(4)

        i += 1

    # Unclosed code block
    if code_lines:
        if not first:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        para.space_before = Pt(6)
        para.space_after = Pt(6)
        for cl in code_lines:
            _run(para, cl + "\n", 11, c["body"], font_name="Consolas")


# ---------------------------------------------------------------------------
# Public entry
# ---------------------------------------------------------------------------

async def generate(slides_content: list[dict], template_id: str, task_id: str) -> str:
    """
    Generate PPT from structured content.

    slides_content: [{"title": "...", "content": "markdown", "notes": "..."}]
    template_id:    theme key
    task_id:        output filename
    Returns:        .pptx file path
    """
    theme = _theme(template_id)
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.333)

    if not slides_content:
        slides_content = [{"title": "内容", "content": "无内容", "notes": ""}]

    blank = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides_content):
        slide = prs.slides.add_slide(blank)
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        notes = slide_data.get("notes", "")

        if i == 0:
            _build_title_slide(slide, title, content[:300], theme)
        else:
            _build_content_slide(slide, title, content, i + 1, theme)

        try:
            ns = slide.notes_slide
            ns.notes_text_frame.text = notes
        except Exception:
            pass

    path = DATA_DIR / "ppts" / f"{task_id}.pptx"
    prs.save(str(path))
    logger.info(f"[PPT] theme={template_id} -> {path}")
    return str(path)
